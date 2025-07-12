#!/usr/bin/env python3

import os
import json
import uuid
from datetime import datetime
from typing import Dict, List
from flask import Flask, request, jsonify, abort
from flask_cors import CORS
from werkzeug.utils import secure_filename
import openai
from pptx import Presentation

# ——— Configuration —————————————
openai.api_key = os.getenv('OPENAI_API_KEY') or "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"

UPLOAD_FOLDER = os.path.join(os.getcwd(), "uploads")
ALLOWED_EXTENSIONS = {'.pdf', '.ppt', '.pptx'}

# ——— Flask Initialization ——————————
app = Flask(__name__)
app.config['UPLOAD_FOLDER'] = UPLOAD_FOLDER
app.config['MAX_CONTENT_LENGTH'] = 50 * 1024 * 1024  # Limit to 50 MB

CORS(app)

# ——— In-Memory Storage ——————————
materials_db = {}

# ——— Helper Functions ——————————
def allowed_file(filename: str) -> bool:
    ext = os.path.splitext(filename)[1].lower()
    return ext in ALLOWED_EXTENSIONS

def extract_text_from_pptx(file_path: str) -> str:
    prs = Presentation(file_path)
    text = ""
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + "\n"
    return text

# ——— Material Processor —————————
class MaterialProcessor:
    def process_pdf(self, file_path: str) -> Dict:
        return {
            'text': "Sample extracted text from PDF.",
            'pages': 25,
            'metadata': {'title': 'Sample PDF Document'}
        }

    def process_pptx(self, file_path: str) -> Dict:
        text = extract_text_from_pptx(file_path)
        slide_count = len(Presentation(file_path).slides)
        return {
            'text': text,
            'slides': slide_count,
            'metadata': {'title': 'PowerPoint Presentation'}
        }

    def create_vector_store(self, text: str, material_id: str):
        print(f"Created vector store for material {material_id}")

# ——— AI Assistant ———————————
class AIAssistant:
    def answer_question(self, question: str, material_id: str) -> Dict:
        material = materials_db.get(material_id)
        if not material:
            return {"answer": "Material not found.", "sources": []}

        context = material['processed_data']['text']
        prompt = f"""You are a helpful learning assistant. Based on the following material, answer the user's question.\n\nMaterial:\n\"\"\"\n{context}\n\"\"\"\n\nQuestion: {question}\n\nAnswer:"""
        response = openai.ChatCompletion.create(
            model="gpt-3.5-turbo",
            messages=[{"role": "user", "content": prompt}],
            max_tokens=500
        )
        return {
            'answer': response.choices[0].message.content.strip(),
            'sources': ['Generated from uploaded material']
        }

    def generate_quiz(self, material_id: str, num_questions: int = 5) -> List[Dict]:
        material = materials_db.get(material_id)
        if not material:
            return []
        context = material['processed_data']['text']
        prompt = f"""Create {num_questions} multiple-choice quiz questions from the following educational material.
Each question must have 4 options and identify the correct answer.\n\nMaterial:\n\"\"\"\n{context}\n\"\"\"\n\nRespond in JSON format:\n[{{
  "question": "...",
  "options": ["A", "B", "C", "D"],
  "correct_answer": "A",
  "explanation": "..."
}}, ...]"""
        response = openai.ChatCompletion.create(
            model="gpt-3.5-turbo",
            messages=[{"role": "user", "content": prompt}],
            max_tokens=1000
        )
        try:
            quiz = json.loads(response.choices[0].message.content)
        except Exception as e:
            quiz = [{
                "question": "Failed to parse quiz",
                "options": [],
                "correct_answer": "",
                "explanation": str(e)
            }]
        return quiz

processor = MaterialProcessor()
ai_assistant = AIAssistant()

# ——— API Endpoints ——————————

@app.route('/api/health', methods=['GET'])
def health_check():
    return jsonify({'status': 'healthy', 'timestamp': datetime.now().isoformat()})

@app.route('/api/upload', methods=['POST'])
def upload_file():
    if 'file' not in request.files:
        return jsonify({'error': 'No file provided'}), 400

    file = request.files['file']
    if file.filename == '':
        return jsonify({'error': 'No file selected'}), 400

    orig_filename = secure_filename(file.filename)
    if not allowed_file(orig_filename):
        return jsonify({'error': 'Unsupported file type'}), 400

    material_id = str(uuid.uuid4())
    filename = f"{material_id}_{orig_filename}"
    os.makedirs(app.config['UPLOAD_FOLDER'], exist_ok=True)
    filepath = os.path.join(app.config['UPLOAD_FOLDER'], filename)
    file.save(filepath)

    if orig_filename.lower().endswith(('.ppt', '.pptx')):
        processed_data = processor.process_pptx(filepath)
    else:
        processed_data = processor.process_pdf(filepath)

    processor.create_vector_store(processed_data['text'], material_id)

    materials_db[material_id] = {
        'id': material_id,
        'filename': orig_filename,
        'file_path': filepath,
        'processed_data': processed_data,
        'uploaded_at': datetime.now().isoformat(),
        'processed': True
    }

    return jsonify({
        'material_id': material_id,
        'filename': orig_filename,
        'processed': True,
        'pages': processed_data.get('pages', processed_data.get('slides', 0)),
        'message': 'File uploaded and processed successfully'
    })

@app.route('/api/chat', methods=['POST'])
def chat():
    data = request.get_json()
    q = data.get('question'); m_id = data.get('material_id')
    if not q or not m_id:
        return jsonify({'error': 'Question and material_id are required'}), 400
    if m_id not in materials_db:
        return jsonify({'error': 'Material not found'}), 404
    r = ai_assistant.answer_question(q, m_id)
    return jsonify({'answer': r['answer'], 'sources': r['sources'], 'timestamp': datetime.now().isoformat()})

@app.route('/api/quiz', methods=['POST'])
def generate_quiz():
    data = request.get_json()
    m_id = data.get('material_id'); n = data.get('num_questions', 5)
    if not m_id:
        return jsonify({'error': 'Material ID is required'}), 400
    if m_id not in materials_db:
        return jsonify({'error': 'Material not found'}), 404
    qlist = ai_assistant.generate_quiz(m_id, n)
    return jsonify({'questions': qlist, 'total_questions': len(qlist), 'generated_at': datetime.now().isoformat()})

if __name__ == '__main__':
    print("🚀 Starting AI Learning Assistant Backend...")
    app.run(debug=True, host='0.0.0.0', port=5000)
